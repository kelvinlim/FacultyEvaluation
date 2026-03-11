"""Build a PowerPoint presentation with one generated PNG table per slide."""

from pathlib import Path
import struct

from pptx import Presentation

from criteria_data import load_tables

PNG_DIR = Path(__file__).with_name("png")
OUTPUT_FILE = Path(__file__).with_name("Faculty_Performance_Criteria_PNGs.pptx")

SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5
MARGIN_INCHES = 0.25
EMU_PER_INCH = 914400


def png_dimensions(path: Path) -> tuple[int, int]:
    with path.open("rb") as file_obj:
        signature = file_obj.read(8)
        if signature != b"\x89PNG\r\n\x1a\n":
            raise ValueError(f"Not a PNG file: {path.name}")
        chunk_length = struct.unpack(">I", file_obj.read(4))[0]
        chunk_type = file_obj.read(4)
        if chunk_type != b"IHDR" or chunk_length < 8:
            raise ValueError(f"Invalid PNG header: {path.name}")
        width, height = struct.unpack(">II", file_obj.read(8))
    return width, height


def inches_to_emu(value: float) -> int:
    return int(value * EMU_PER_INCH)


def main() -> None:
    tables = load_tables()
    image_paths = []

    for mission_name in tables:
        image_name = mission_name.replace(" ", "_") + ".png"
        image_path = PNG_DIR / image_name
        if not image_path.exists():
            raise FileNotFoundError(f"Missing PNG for slide generation: {image_path}")
        image_paths.append(image_path)

    presentation = Presentation()
    presentation.slide_width = inches_to_emu(SLIDE_WIDTH_INCHES)
    presentation.slide_height = inches_to_emu(SLIDE_HEIGHT_INCHES)
    blank_layout = presentation.slide_layouts[6]

    usable_width = presentation.slide_width - (2 * inches_to_emu(MARGIN_INCHES))
    usable_height = presentation.slide_height - (2 * inches_to_emu(MARGIN_INCHES))

    for image_path in image_paths:
        slide = presentation.slides.add_slide(blank_layout)
        image_width_px, image_height_px = png_dimensions(image_path)
        width_scale = usable_width / image_width_px
        height_scale = usable_height / image_height_px
        scale = min(width_scale, height_scale)

        image_width = int(image_width_px * scale)
        image_height = int(image_height_px * scale)
        left = int((presentation.slide_width - image_width) / 2)
        top = int((presentation.slide_height - image_height) / 2)

        slide.shapes.add_picture(str(image_path), left, top, width=image_width, height=image_height)

    presentation.save(OUTPUT_FILE)
    print(f"Done! Saved to {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
